"""Build the AI Ninja Program executive deck PPTX from EXEC_SUMMARY_DECK.md content.

Reuses the visual template from the lesson 1 lecture deck so the exec deck inherits
the same brand layouts (colors, fonts, master slides).
"""

import copy
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parent.parent
TEMPLATE = ROOT / "curriculum" / "stage1_classic_ml" / "01_what_is_ml" / "Lecture-1-What-is-ML.pptx"
OUTPUT = ROOT / "docs" / "AI-Ninja-Program-Exec-Deck.pptx"

# Brand palette (mirrors portal CSS vars)
CYAN   = RGBColor(0x06, 0xD6, 0xE0)
VIOLET = RGBColor(0xA8, 0x55, 0xF7)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
DARK   = RGBColor(0x1F, 0x29, 0x37)
GREY   = RGBColor(0x4B, 0x55, 0x63)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def remove_all_slides(prs):
    """Strip every slide from the loaded template, leaving layouts/masters intact."""
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.attrib["{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"]
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def layout_by_name(prs, name):
    for layout in prs.slide_layouts:
        if layout.name.strip() == name.strip():
            return layout
    raise KeyError(f"Layout {name!r} not found. Available: {[l.name for l in prs.slide_layouts]}")


def get_placeholder(slide, idx):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    raise KeyError(f"Placeholder idx={idx} not found on slide. Have: {[p.placeholder_format.idx for p in slide.placeholders]}")


def set_text(placeholder, text):
    """Replace placeholder text with a single line, preserving template formatting."""
    tf = placeholder.text_frame
    tf.clear()
    tf.paragraphs[0].text = text


def set_bullets(placeholder, bullets):
    """Populate a body placeholder with bullets. Bullets is a list of (text, level) or strings."""
    tf = placeholder.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = level


def add_speaker_notes(slide, text):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = text


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------

def build():
    prs = Presentation(str(TEMPLATE))
    remove_all_slides(prs)

    layouts = {layout.name.strip(): layout for layout in prs.slide_layouts}
    print("Available layouts:", list(layouts.keys()))

    # ----- SLIDE 1: Title -----
    slide = prs.slides.add_slide(layouts["Title Slide"])
    set_text(get_placeholder(slide, 14), "AI Ninja Program")
    set_text(get_placeholder(slide, 13), "Building deep AI/ML capability across our Security Engineering team")
    set_text(get_placeholder(slide, 1), "Khalid Alshawwaf  ·  For: Head of Security Engineering")
    try:
        set_text(get_placeholder(slide, 15), "~7 minute proposal + Q&A")
    except KeyError:
        pass
    add_speaker_notes(
        slide,
        "This is a proposal for an internal training program. I'll walk you through the "
        "problem, the program, why I think it will work, and what I need from you."
    )

    # ----- SLIDE 2: The problem -----
    slide = prs.slides.add_slide(layouts["Standard Title and Bullets"])
    set_text(get_placeholder(slide, 0), "Every customer conversation now has AI in it")
    set_bullets(get_placeholder(slide, 1), [
        "Vendors pitching \"AI-powered EDR\" and \"ML-driven threat detection\"",
        "CISOs asking us to evaluate their LLM and RAG use cases",
        "Competitors claiming their anomaly detection is \"machine learning\"",
        "Customers expecting us to engage on AI as peers, not observers",
        "",
        "Our engineers are world-class at firewalls, SIEM, and EDR.",
        "The AI layer is increasingly the part that decides who wins the deal.",
    ])
    add_speaker_notes(
        slide,
        "The gap is not awareness — everybody knows AI is important. The gap is technical depth. "
        "Reading a blog about transformers does not let you push back when a vendor says 'our model "
        "catches zero-days' and you need to ask which model, trained on what data, with what "
        "false-positive rate."
    )

    # ----- SLIDE 3: The proposal -----
    slide = prs.slides.add_slide(layouts["Standard Title and Bullets"])
    set_text(get_placeholder(slide, 0), "A 15-week internal training program")
    set_bullets(get_placeholder(slide, 1), [
        "Takes a security engineer with zero ML background and turns them into someone who can:",
        ("Explain  AI/ML to a customer in plain language", 1),
        ("Evaluate  vendor AI claims with technical credibility", 1),
        ("Build  working ML classifiers for security use cases", 1),
        ("Critique  neural network architectures", 1),
        ("Demo  a live AI-powered security assistant", 1),
        ("Position  Check Point AI Security products with technical authority", 1),
        "",
        "Operating principle: people learn AI by manipulating it, not reading about it.",
    ])
    add_speaker_notes(
        slide,
        "The whole curriculum is already built and tested. 21 lessons, hundreds of pages of "
        "hands-on labs, real datasets — not slides. I can demo the portal after this."
    )

    # ----- SLIDE 4: Structure (Title Only + custom table) -----
    slide = prs.slides.add_slide(layouts["Title Only"])
    set_text(get_placeholder(slide, 0), "21 lessons  ·  5 stages  ·  3 progressive certification tiers")

    # Build a table manually
    rows, cols = 4, 3
    left = Inches(0.7)
    top  = Inches(2.0)
    width  = Inches(11.6)
    height = Inches(2.6)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    table.columns[0].width = Inches(2.4)
    table.columns[1].width = Inches(1.6)
    table.columns[2].width = Inches(7.6)

    headers = ["Tier", "Weeks", "What graduates can do"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.bold = True
                run.font.size = Pt(16)
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = VIOLET

    rows_data = [
        ("Tier 1 — AI Foundations",  "1–7",   "Hold credible AI conversations with customers; evaluate vendor claims"),
        ("Tier 2 — AI Practitioner", "8–10",  "Prototype neural networks; assess AI solution architectures"),
        ("Tier 3 — AI Ninja",        "11–15", "Build AI-powered security tools; lead AI engagements end-to-end"),
    ]
    for r, row in enumerate(rows_data, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(14)
                    run.font.color.rgb = DARK
            if c == 0:
                for p in cell.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.bold = True

    # Caption text below table
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(4.9), Inches(11.6), Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "5 stages: AI Positioning → Classic ML → Intermediate ML → Neural Networks → Generative AI / RAG → Applied Check Point AI"
    for run in p.runs:
        run.font.size = Pt(14)
        run.font.italic = True
        run.font.color.rgb = GREY

    add_speaker_notes(
        slide,
        "The tier system matters because not everyone needs the full depth. A senior architect "
        "who only needs sales credibility can stop at Tier 1 in 7 weeks. Engineers building tools "
        "go all the way to Tier 3. We get more participation because we respect their time."
    )

    # ----- SLIDE 5: Why it will work -----
    slide = prs.slides.add_slide(layouts["Standard Title and Bullets"])
    set_text(get_placeholder(slide, 0), "Three things that separate this from \"watch a Coursera playlist\"")
    set_bullets(get_placeholder(slide, 1), [
        "1. Graded assessment gates between stages",
        ("Quiz · mini-project · architecture review · live capstone demo", 1),
        ("People can't just run solution files and graduate — they have to demonstrate understanding", 1),
        "2. Real security data, not toy datasets",
        ("Real CVEs · real phishing URLs · MITRE ATT&CK · packet captures", 1),
        ("Capstones use data graduates could demo to a customer the next day", 1),
        "3. Every Ninja graduate ships a portable demo",
        ("Laptop-ready RAG security assistant + 5-minute customer-facing script", 1),
        ("The deliverable is not a certificate — it's a sales tool", 1),
    ])
    add_speaker_notes(
        slide,
        "The capstone is the part I'm most excited about. Imagine an engineer walking into a "
        "customer meeting, pulling out their laptop, and saying 'let me show you what I built in "
        "our AI program.' That's worth 10x a certificate."
    )

    # ----- SLIDE 6: The ask (Title Only + table) -----
    slide = prs.slides.add_slide(layouts["Title Only"])
    set_text(get_placeholder(slide, 0), "Four things I'm asking for")

    rows, cols = 5, 3
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.7), Inches(1.9), Inches(11.6), Inches(3.6))
    table = table_shape.table
    table.columns[0].width = Inches(0.6)
    table.columns[1].width = Inches(5.2)
    table.columns[2].width = Inches(5.8)

    headers = ["#", "Ask", "What it costs"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.bold = True
                run.font.size = Pt(16)
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = CYAN

    asks = [
        ("1", "Endorsement to recruit a pilot cohort",
              "Your name on the program announcement"),
        ("2", "Time allocation: 5–8 hrs/week per participant for 15 weeks",
              "~12–16 engineers across the Americas"),
        ("3", "Sponsorship so participation reads as career investment",
              "One internal email + visible support"),
        ("4", "Review panel seat for the final capstone presentations",
              "~2 hours at week 15, you + 2–3 stakeholders"),
    ]
    for r, row in enumerate(asks, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(14)
                    run.font.color.rgb = DARK
            if c == 0:
                for p in cell.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.CENTER
                    for run in p.runs:
                        run.font.bold = True

    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(5.8), Inches(11.6), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "No external budget required. Curriculum is built. Portal is built. Materials are tested."
    for run in p.runs:
        run.font.size = Pt(15)
        run.font.bold = True
        run.font.color.rgb = VIOLET

    add_speaker_notes(
        slide,
        "I'm not asking for a budget line. The curriculum already exists, the portal already runs, "
        "the labs are already tested. What I need from you is air cover and access to the right engineers."
    )

    # ----- SLIDE 7: Success -----
    slide = prs.slides.add_slide(layouts["Summary"])
    set_text(get_placeholder(slide, 0), "What success looks like in 15 weeks")
    set_bullets(get_placeholder(slide, 12), [
        "12–16 certified engineers who can hold technical AI conversations with any customer in the Americas",
        "Each graduate carries a working AI security demo on their laptop, ready for customer meetings",
        "A repeatable curriculum ready for the next cohort with minimal setup",
        "Clear visibility into Tier 1 (sales-ready), Tier 2 (architecture-ready), Tier 3 (build-ready) capability",
        "",
        "The customers we sell to are judging us on AI fluency. This is the lowest-risk, highest-leverage way to build it.",
    ])
    add_speaker_notes(
        slide,
        "I'd love 20 minutes after this to walk you through the actual program portal so you can "
        "see what graduates will be working with. Happy to take any questions now."
    )

    # ----- SLIDE 8: Close -----
    slide = prs.slides.add_slide(layouts["Close"])
    set_text(get_placeholder(slide, 14), "Want to see the portal?")
    set_text(get_placeholder(slide, 13), "20 minutes, in person — I'll walk you through the labs and the capstone deliverable.")

    # Save
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Wrote {OUTPUT}")
    print(f"Slide count: {len(prs.slides)}")


if __name__ == "__main__":
    build()
